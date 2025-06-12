import streamlit as st
import PyPDF2
from pptx import Presentation
import google.generativeai as genai
import re
import json
import os
from datetime import datetime
import hashlib
import base64
from supabase import create_client, Client

# --- Supabase Setup ---
SUPABASE_URL = st.secrets["SUPABASE_URL"]
SUPABASE_KEY = st.secrets["SUPABASE_KEY"]

@st.cache_resource
def get_supabase_client():
    return create_client(SUPABASE_URL, SUPABASE_KEY)

def save_quiz_to_supabase(quiz_id, questions):
    supabase: Client = get_supabase_client()
    # Upsert (insert or update) the quiz
    supabase.table("quizzes").upsert({"quiz_id": quiz_id, "questions_json": questions}).execute()

def load_quiz_from_supabase(quiz_id):
    supabase: Client = get_supabase_client()
    result = supabase.table("quizzes").select("questions_json").eq("quiz_id", quiz_id).execute()
    if result.data and len(result.data) > 0:
        return result.data[0]["questions_json"]
    return None

# Get API key from Streamlit secrets
GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]

# Create a directory for storing quizzes if it doesn't exist
if not os.path.exists('quizzes'):
    os.makedirs('quizzes')

def get_db_connection():
    conn = sqlite3.connect("quiz.db", check_same_thread=False)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS quizzes (
            quiz_id TEXT PRIMARY KEY,
            questions_json TEXT
        )
    """)
    return conn

def save_quiz_to_file(quiz_id, questions):
    """Save quiz to a JSON file"""
    file_path = f'quizzes/{quiz_id}.json'
    with open(file_path, 'w', encoding='utf-8') as f:
        json.dump(questions, f, ensure_ascii=False, indent=2)

def load_quiz_from_file(quiz_id):
    """Load quiz from a JSON file"""
    file_path = f'quizzes/{quiz_id}.json'
    if os.path.exists(file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def extract_text_from_pdf(pdf_file):
    text = ""
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + " "
    return text

def extract_text_from_pptx(pptx_file):
    text = ""
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + " "
    return text

def clean_text(text, max_words=600):
    lines = list(dict.fromkeys(text.splitlines()))
    text = " ".join(lines)
    words = text.split()
    if len(words) > max_words:
        text = " ".join(words[:max_words])
    return text

def generate_mcqs_with_gemini(text, num_questions):
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = (
        f"You are a cybersecurity teacher. Based on the following content, generate {num_questions} high-quality, conceptual multiple-choice questions (MCQs) for students. "
        "Each question should have 4 options (A, B, C, D), only one correct answer, and a short explanation. "
        "Do NOT use fill-in-the-blank or word replacement. Make the questions conceptual and relevant to the content. "
        "Format:\n"
        "Q: <question>\nA) <option1>\nB) <option2>\nC) <option3>\nD) <option4>\nAnswer: <A/B/C/D>\nExplanation: <short explanation>\n"
        "Content:\n"
        f"{text}\n"
    )
    response = model.generate_content(prompt)
    return response.text

def parse_mcqs(gemini_response):
    questions = []
    blocks = re.split(r"\nQ: ", "\n" + gemini_response)
    for block in blocks[1:]:
        lines = block.strip().split("\n")
        q = lines[0].strip()
        options = []
        for line in lines[1:5]:
            if line[:2] in ["A)", "B)", "C)", "D)"]:
                options.append(line[3:].strip())
        answer_line = next((l for l in lines if l.startswith("Answer:")), None)
        explanation_line = next((l for l in lines if l.startswith("Explanation:")), None)
        if answer_line and explanation_line and len(options) == 4:
            answer_letter = answer_line.split(":", 1)[1].strip()
            # Only accept A/B/C/D as valid answers
            if answer_letter in "ABCD":
                correct_index = "ABCD".index(answer_letter)
            else:
                # If not valid, skip this question
                continue
            explanation = explanation_line.split(":", 1)[1].strip()
            questions.append({
                "question": q,
                "options": options,
                "correct_answer": correct_index,
                "explanation": explanation
            })
    return questions

def generate_quiz_id(questions):
    quiz_str = json.dumps(questions, sort_keys=True)
    return base64.urlsafe_b64encode(hashlib.md5(quiz_str.encode()).digest()).decode()[:8]

def show_quiz_interface(questions):
    if 'current_question' not in st.session_state:
        st.session_state.current_question = 0
        st.session_state.score = 0
        st.session_state.checked = False
        st.session_state.selected_option = None
        st.session_state.explanation_shown = False
        st.session_state.answers = [None] * len(questions)
        st.session_state.attempted = [False] * len(questions)

    q_idx = st.session_state.current_question
    q = questions[q_idx]
    total_questions = len(questions)

    st.progress((q_idx+1)/total_questions, text=f"{q_idx+1}/{total_questions}")
    st.subheader(f"Question {q_idx+1}")
    st.write(q['question'])

    answer_locked = st.session_state.answers[q_idx] is not None
    selected_option = st.radio(
        "Choose your answer:",
        q['options'],
        index=st.session_state.selected_option if st.session_state.selected_option is not None else 0,
        key=f"radio_{q_idx}",
        disabled=answer_locked
    )

    col1, col2 = st.columns([6,1])
    with col2:
        submit_btn = st.button("Check", key=f"check_{q_idx}", disabled=answer_locked)

    if submit_btn and not answer_locked:
        st.session_state.checked = True
        st.session_state.selected_option = q['options'].index(selected_option)
        st.session_state.answers[q_idx] = st.session_state.selected_option
        st.session_state.attempted[q_idx] = True
        st.session_state.explanation_shown = False

    if st.session_state.answers[q_idx] is not None:
        selected = st.session_state.answers[q_idx]
        if selected == q['correct_answer']:
            st.success("Correct!")
            if not st.session_state.get(f'scored_{q_idx}', False):
                st.session_state.score += 1
                st.session_state[f'scored_{q_idx}'] = True
        else:
            st.error("Incorrect!")

        # Improved button layout: Explanation and Next on the same line, centered
        exp_col, next_col, _ = st.columns([2,2,6])
        with exp_col:
            show_exp = st.button("💡 Explanation", key=f"explanation_{q_idx}")
        with next_col:
            next_btn = st.button("➡️ Next", key=f"next_{q_idx}")
        if show_exp:
            st.session_state.explanation_shown = True
        if st.session_state.explanation_shown:
            st.info(f"Explanation: {q['explanation']}")
        if next_btn:
            st.session_state.current_question += 1
            st.session_state.checked = False
            st.session_state.selected_option = None
            st.session_state.explanation_shown = False
            if st.session_state.current_question >= total_questions:
                show_result_summary(questions)
                st.session_state.clear()

def show_result_summary(questions):
    total = len(questions)
    attempted = sum(1 for a in st.session_state.answers if a is not None)
    correct = st.session_state.score
    wrong = attempted - correct
    percent = (correct / total) * 100 if total > 0 else 0
    st.balloons()
    st.success(f"Quiz completed!")
    st.markdown(f"**Total Questions:** {total}")
    st.markdown(f"**Attempted:** {attempted}")
    st.markdown(f"**Correct:** {correct}")
    st.markdown(f"**Wrong:** {wrong}")
    st.markdown(f"**Score:** {percent:.2f}%")

def main():
    st.set_page_config(page_title="Cybersecurity Quiz Generator", layout="centered")
    st.markdown("""
        <div style="display: flex; justify-content: space-between; align-items: flex-end;">
            <div>
                <span style="font-size: 2.2rem; font-weight: 700;">Cybersecurity Quiz Generator</span>
                <span style="font-size: 1.1rem; font-weight: 400; color: #FF4B4B; vertical-align: super;">(AI-Powered)</span>
            </div>
            <div style="text-align: right; color: #888; font-size: 17px; margin-bottom: 10px; min-width: 320px;">
                developed by : <b>Badri Tamang</b> {S+, SCS-C02}
            </div>
        </div>
    """, unsafe_allow_html=True)

    query_params = st.query_params
    quiz_id = query_params.get("quiz", [None])[0] if "quiz" in query_params else None

    if quiz_id:
        questions = load_quiz_from_supabase(quiz_id)
        if questions:
            show_quiz_interface(questions)
        else:
            st.error(
                "Quiz not found. Please generate a new quiz.\n\n"
                "Note: Quizzes are only available if generated and saved in Supabase. "
                "If you just created a quiz, please share the new link."
            )
            quiz_id = None

    if not quiz_id:
        st.write("Welcome! Upload a file to generate a new quiz.")
        uploaded_file = st.file_uploader("Choose a file", type=['pptx', 'pdf'])
        if uploaded_file:
            num_questions = st.slider("Number of questions to generate", min_value=5, max_value=20, value=10)
            if st.button("Generate Quiz"):
                with st.spinner("Extracting text and generating questions. Please wait..."):
                    if uploaded_file.name.endswith('.pdf'):
                        text = extract_text_from_pdf(uploaded_file)
                    else:
                        text = extract_text_from_pptx(uploaded_file)
                    text = clean_text(text)
                    gemini_response = generate_mcqs_with_gemini(text, num_questions)
                    questions = parse_mcqs(gemini_response)
                    if not questions:
                        st.error("Could not generate questions. Try with a different file or fewer questions.")
                    else:
                        quiz_id = generate_quiz_id(questions)
                        save_quiz_to_supabase(quiz_id, questions)
                        st.success("Quiz generated successfully!")
                        st.markdown("### To share this quiz with students:")
                        st.markdown(f"1. Your quiz is now saved and can be accessed using this URL:\n   `https://cyberquizbt.streamlit.app/?quiz={quiz_id}`")
                        st.markdown(f"**Quiz ID:** `{quiz_id}`")
                        st.markdown("Share this complete URL with your students.")

if __name__ == "__main__":
    main() 